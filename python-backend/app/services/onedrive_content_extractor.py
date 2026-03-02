"""OneDrive file content extraction service.

Extracts text from various file formats downloaded from OneDrive.
"""

import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)

# Limits to prevent memory explosion
MAX_EXTRACTION_BYTES = 30 * 1024 * 1024  # 30 MB
MAX_XLSX_ROWS_PER_SHEET = 10_000
MAX_PDF_PAGES = 500

# Legacy Office formats that cannot be parsed natively
LEGACY_OFFICE_EXTENSIONS = (".doc", ".xls", ".ppt", ".wps", ".wks")


class OneDriveContentExtractor:
    """Extracts text content from OneDrive file bytes."""

    def extract(self, content: bytes, mime_type: str, file_name: str = "") -> dict:
        """Extract text from file content based on MIME type.

        Returns dict with keys: text, char_count, method
        """
        # Size guard to prevent OOM
        if len(content) > MAX_EXTRACTION_BYTES:
            logger.warning(
                "File too large for extraction: %s (%d bytes, limit %d)",
                file_name, len(content), MAX_EXTRACTION_BYTES,
            )
            return {"text": "", "char_count": 0, "method": "too_large"}

        mime_lower = mime_type.lower()
        name_lower = file_name.lower()
        text = ""
        method = "unknown"

        try:
            # Legacy Office formats — not natively parseable
            if name_lower.endswith(LEGACY_OFFICE_EXTENSIONS):
                return {
                    "text": "",
                    "char_count": 0,
                    "method": "legacy_unsupported",
                }

            if "text/" in mime_lower or self._is_text_file(file_name):
                text = content.decode("utf-8", errors="replace")
                method = "text"
            elif "pdf" in mime_lower or name_lower.endswith(".pdf"):
                text = self._extract_pdf(content)
                method = "pdf"
            elif "wordprocessingml" in mime_lower or name_lower.endswith(".docx"):
                text = self._extract_docx(content)
                method = "docx"
            elif "presentationml" in mime_lower or name_lower.endswith(".pptx"):
                text = self._extract_pptx(content)
                method = "pptx"
            elif "spreadsheetml" in mime_lower or name_lower.endswith(".xlsx"):
                text = self._extract_xlsx(content)
                method = "xlsx"
            elif "msword" in mime_lower or "ms-excel" in mime_lower or "ms-powerpoint" in mime_lower:
                # Legacy MIME types (application/msword, application/vnd.ms-excel, etc.)
                return {
                    "text": "",
                    "char_count": 0,
                    "method": "legacy_unsupported",
                }
            else:
                # Try as text fallback
                try:
                    text = content.decode("utf-8", errors="strict")
                    method = "text_fallback"
                except (UnicodeDecodeError, ValueError):
                    text = ""
                    method = "unsupported"
        except Exception as e:
            logger.warning("Content extraction failed for %s: %s", file_name, e)
            text = ""
            method = "error"

        return {"text": text, "char_count": len(text), "method": method}

    def _is_text_file(self, name: str) -> bool:
        return name.lower().endswith((
            ".txt", ".md", ".csv", ".json", ".xml",
            ".yaml", ".yml", ".html", ".htm", ".log",
            ".ini", ".cfg", ".toml", ".rst", ".tex",
        ))

    def _extract_pdf(self, content: bytes) -> str:
        from PyPDF2 import PdfReader

        reader = PdfReader(io.BytesIO(content))
        texts = []
        for i, page in enumerate(reader.pages):
            if i >= MAX_PDF_PAGES:
                texts.append(f"... (truncated at {MAX_PDF_PAGES} pages)")
                break
            texts.append(page.extract_text() or "")
        return "\n".join(texts)

    def _extract_docx(self, content: bytes) -> str:
        from docx import Document

        doc = Document(io.BytesIO(content))
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text for cell in row.cells)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)

    def _extract_pptx(self, content: bytes) -> str:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(content))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
                # Extract table content
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = "\t".join(cell.text for cell in row.cells)
                        if row_text.strip():
                            texts.append(row_text)
        return "\n".join(texts)

    def _extract_xlsx(self, content: bytes) -> str:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        texts = []
        for ws in wb.worksheets:
            texts.append(f"--- Sheet: {ws.title} ---")
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(c) if c is not None else "" for c in row)
                if row_text.strip():
                    texts.append(row_text)
                    row_count += 1
                    if row_count >= MAX_XLSX_ROWS_PER_SHEET:
                        texts.append(f"... (truncated at {MAX_XLSX_ROWS_PER_SHEET} rows)")
                        break
        wb.close()
        return "\n".join(texts)
