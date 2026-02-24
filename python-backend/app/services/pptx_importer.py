"""
PPTX Importer — parses a .pptx file into PresentationSlideContent dicts.

Uses python-pptx to extract shapes (text, image, rect, line, group).
Uploads embedded images to R2 via R2StorageService.upload_bytes.
Returns ImportResult with slides and fidelity_warnings.
"""
import io
from dataclasses import dataclass
from typing import Optional
from uuid import uuid4

import structlog
from pptx import Presentation
from pptx.enum.dml import MSO_COLOR_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE

from app.services.presentation_importer import ImportResult
from app.services.r2_storage_service import R2StorageService

logger = structlog.get_logger(__name__)

# Aspect ratio preset map: round(width/height, 3) → (label, canvas_width, canvas_height)
PRESET_MAP: dict[float, tuple[str, int, int]] = {
    1.778: ("16:9", 1280, 720),
    0.563: ("9:16", 720, 1280),
    1.333: ("4:3", 1024, 768),
    0.75: ("3:4", 768, 1024),
    0.8: ("4:5", 960, 1200),
    1.25: ("5:4", 1250, 1000),
    1.0: ("1:1", 1080, 1080),
}

MAX_PPTX_SIZE = 50_000_000  # 50 MB — guard against zip bombs / memory exhaustion

CONTENT_TYPE_TO_EXT: dict[str, str] = {
    "image/png": "png",
    "image/jpeg": "jpg",
    "image/gif": "gif",
    "image/webp": "webp",
    "image/bmp": "bmp",
    "image/tiff": "tiff",
}


# ---------------------------------------------------------------------------
# Pure coordinate conversion functions (module-level for direct test import)
# ---------------------------------------------------------------------------


def _emu_to_px(emu: float, dpi: int = 96) -> float:
    """Convert EMU (English Metric Units) to pixels at the given DPI.

    1 inch = 914400 EMU = dpi pixels.
    """
    return emu * dpi / 914_400


def _scale_to_canvas(emu: float, slide_emu: float, canvas_px: int) -> int:
    """Scale an EMU dimension proportionally to the canvas pixel space."""
    return round(_emu_to_px(emu) * canvas_px / _emu_to_px(slide_emu))


def _pt_to_canvas_px(pt: float) -> int:
    """Convert font size in points to pixels, clamped to [8, 512]."""
    return max(8, min(512, round(pt * 4 / 3)))


# ---------------------------------------------------------------------------
# Canvas detection
# ---------------------------------------------------------------------------


@dataclass
class _CanvasSize:
    """Resolved canvas dimensions for a presentation."""

    preset: Optional[str]  # e.g. "16:9", or None if unknown
    width: int  # canvas width in px
    height: int  # canvas height in px


def _detect_canvas(slide_width_emu: int, slide_height_emu: int) -> _CanvasSize:
    """Look up slide dimensions in PRESET_MAP; fall back to natural px."""
    ratio = round(slide_width_emu / slide_height_emu, 3)
    if ratio in PRESET_MAP:
        label, w, h = PRESET_MAP[ratio]
        return _CanvasSize(preset=label, width=w, height=h)
    return _CanvasSize(
        preset=None,
        width=round(_emu_to_px(slide_width_emu)),
        height=round(_emu_to_px(slide_height_emu)),
    )


# ---------------------------------------------------------------------------
# Warning cap helper
# ---------------------------------------------------------------------------


def _cap_warnings(warnings: list[str]) -> list[str]:
    """Cap warnings list at 25 items with a summary item if truncated."""
    if len(warnings) <= 25:
        return warnings
    overflow = len(warnings) - 24
    return warnings[:24] + [f"... and {overflow} more warnings"]


# ---------------------------------------------------------------------------
# PptxImporter class
# ---------------------------------------------------------------------------


class PptxImporter:
    """Parse a PPTX file into a list of PresentationSlideContent dicts.

    Args:
        r2_service: R2StorageService singleton for uploading embedded images.
    """

    def __init__(self, r2_service: R2StorageService) -> None:
        self._r2 = r2_service

    async def import_file(self, pptx_bytes: bytes, s3_prefix: str) -> ImportResult:
        """Parse pptx_bytes and return an ImportResult.

        Args:
            pptx_bytes: Raw bytes of the .pptx file.
            s3_prefix: Prefix for S3 keys, e.g.
                       "{tenant_id}/presentations/imports/{conversion_id}".

        Raises:
            ImportError: If pptx_bytes is not a valid .pptx file.
        """
        if len(pptx_bytes) > MAX_PPTX_SIZE:
            raise ImportError(
                f"File too large: {len(pptx_bytes):,} bytes (limit: {MAX_PPTX_SIZE:,} bytes)"
            )
        try:
            prs = Presentation(io.BytesIO(pptx_bytes))
        except (Exception, MemoryError):
            raise ImportError("The uploaded file is not a valid .pptx file")

        slide_width_emu = int(prs.slide_width)
        slide_height_emu = int(prs.slide_height)
        canvas = _detect_canvas(slide_width_emu, slide_height_emu)

        all_slides: list[dict] = []
        all_warnings: list[str] = []

        for slide_num, slide in enumerate(prs.slides, start=1):
            elements, warnings = await self._parse_slide(
                slide, canvas, slide_width_emu, slide_height_emu, s3_prefix, slide_num
            )
            all_warnings.extend(warnings)
            all_slides.append(
                {
                    "canvasPreset": canvas.preset,
                    "canvasWidth": canvas.width,
                    "canvasHeight": canvas.height,
                    "elements": elements,
                }
            )

        return ImportResult(
            slides=all_slides,
            fidelity_warnings=_cap_warnings(all_warnings),
        )

    async def _parse_slide(
        self,
        slide,
        canvas: _CanvasSize,
        slide_width_emu: int,
        slide_height_emu: int,
        s3_prefix: str,
        slide_num: int,
    ) -> tuple[list[dict], list[str]]:
        """Parse all shapes on one slide. Returns (elements, warnings)."""
        return await self._parse_shapes(
            slide.shapes,
            parent_left_emu=0,
            parent_top_emu=0,
            canvas=canvas,
            slide_width_emu=slide_width_emu,
            slide_height_emu=slide_height_emu,
            s3_prefix=s3_prefix,
            slide_num=slide_num,
        )

    async def _parse_shapes(
        self,
        shapes,
        parent_left_emu: int,
        parent_top_emu: int,
        canvas: _CanvasSize,
        slide_width_emu: int,
        slide_height_emu: int,
        s3_prefix: str,
        slide_num: int,
    ) -> tuple[list[dict], list[str]]:
        """Recursively parse a shape collection with accumulated parent offsets."""
        elements: list[dict] = []
        warnings: list[str] = []

        for shape in shapes:
            shape_type = shape.shape_type

            abs_left = (shape.left or 0) + parent_left_emu
            abs_top = (shape.top or 0) + parent_top_emu
            x = _scale_to_canvas(abs_left, slide_width_emu, canvas.width)
            y = _scale_to_canvas(abs_top, slide_height_emu, canvas.height)
            w = _scale_to_canvas(shape.width or 0, slide_width_emu, canvas.width)
            h = _scale_to_canvas(shape.height or 0, slide_height_emu, canvas.height)

            if shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                elem = self._parse_textbox(shape, x, y, w, h)
                if elem:
                    elements.append(elem)

            elif shape_type == MSO_SHAPE_TYPE.PICTURE:
                elem, warn = await self._parse_picture(shape, x, y, w, h, s3_prefix, slide_num)
                if elem:
                    elements.append(elem)
                warnings.extend(warn)

            elif shape_type == MSO_SHAPE_TYPE.LINKED_PICTURE:
                warnings.append(
                    f"Slide {slide_num}: Linked image skipped (external link, blob not embedded)"
                )

            elif shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                elems, warns = self._parse_auto_shape(shape, x, y, w, h, slide_num)
                elements.extend(elems)
                warnings.extend(warns)

            elif shape_type == MSO_SHAPE_TYPE.LINE:
                elem = self._parse_line(shape, x, y, w, h)
                if elem:
                    elements.append(elem)

            elif shape_type == MSO_SHAPE_TYPE.GROUP:
                group_left = shape.left or 0
                group_top = shape.top or 0
                child_elems, child_warns = await self._parse_shapes(
                    shape.shapes,
                    parent_left_emu + group_left,
                    parent_top_emu + group_top,
                    canvas,
                    slide_width_emu,
                    slide_height_emu,
                    s3_prefix,
                    slide_num,
                )
                elements.extend(child_elems)
                warnings.extend(child_warns)

            elif shape_type == MSO_SHAPE_TYPE.TABLE:
                warnings.append(f"Slide {slide_num}: Table dropped (not supported)")

            elif shape_type == MSO_SHAPE_TYPE.CHART:
                warnings.append(f"Slide {slide_num}: Chart dropped (not supported)")

            else:
                warnings.append(
                    f"Slide {slide_num}: Unknown shape type {shape_type} skipped"
                )

        return elements, warnings

    def _parse_textbox(self, shape, x: int, y: int, w: int, h: int) -> Optional[dict]:
        """Parse a TEXT_BOX shape into a text element dict."""
        try:
            tf = shape.text_frame
        except Exception:
            return None

        paragraphs = tf.paragraphs
        lines: list[str] = []
        for para in paragraphs:
            para_text = "".join(run.text for run in para.runs)
            lines.append(para_text)

        text = "\n".join(lines)[:10_000]
        if not text.strip():
            return None

        style: dict = {}
        # Find first non-empty run for style extraction
        for para in paragraphs:
            for run in para.runs:
                if run.text.strip():
                    font = run.font
                    try:
                        if (
                            font.color is not None
                            and font.color.type == MSO_COLOR_TYPE.RGB
                        ):
                            style["color"] = f"#{font.color.rgb}"
                    except Exception:
                        pass
                    try:
                        if font.size is not None:
                            style["fontSize"] = _pt_to_canvas_px(font.size.pt)
                    except Exception:
                        pass
                    try:
                        if font.bold is not None:
                            style["bold"] = font.bold
                    except Exception:
                        pass
                    try:
                        if font.italic is not None:
                            style["italic"] = font.italic
                    except Exception:
                        pass
                    break
            else:
                continue
            break

        return {
            "type": "text",
            "x": x,
            "y": y,
            "width": w,
            "height": h,
            "content": text,
            "style": style,
        }

    async def _parse_picture(
        self, shape, x: int, y: int, w: int, h: int, s3_prefix: str, slide_num: int
    ) -> tuple[Optional[dict], list[str]]:
        """Parse a PICTURE shape: upload blob to R2, return image element."""
        try:
            blob = shape.image.blob
            content_type = shape.image.content_type or "image/png"
            ext = CONTENT_TYPE_TO_EXT.get(content_type, "png")
            key = f"{s3_prefix}/images/{uuid4()}.{ext}"
            url = await self._r2.upload_bytes(key, blob, content_type)
            return {
                "type": "image",
                "x": x,
                "y": y,
                "width": w,
                "height": h,
                "src": url,
            }, []
        except Exception as e:
            logger.warning("pptx_picture_parse_error", error=str(e), slide=slide_num)
            return None, [f"Slide {slide_num}: Image could not be extracted ({e})"]

    def _parse_auto_shape(
        self, shape, x: int, y: int, w: int, h: int, slide_num: int
    ) -> tuple[list[dict], list[str]]:
        """Parse an AUTO_SHAPE: returns rect/oval element and any warnings."""
        elems: list[dict] = []
        warns: list[str] = []

        try:
            ast = shape.auto_shape_type
        except Exception:
            warns.append(f"Slide {slide_num}: Shape type could not be determined, skipped")
            return elems, warns

        is_rect = ast in (MSO_AUTO_SHAPE_TYPE.RECTANGLE, MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE)
        is_oval = ast == MSO_AUTO_SHAPE_TYPE.OVAL

        if not is_rect and not is_oval:
            warns.append(f"Slide {slide_num}: Shape type '{ast}' not supported")
            return elems, warns

        if is_oval:
            warns.append(f"Slide {slide_num}: Oval approximated as rectangle")

        # Extract fill color
        fill = "#cccccc"
        try:
            if hasattr(shape, "fill"):
                fill_obj = shape.fill
                if fill_obj.type is not None:
                    try:
                        fc = fill_obj.fore_color
                        if fc.type == MSO_COLOR_TYPE.RGB:
                            fill = f"#{fc.rgb}"
                    except Exception:
                        pass
        except Exception:
            pass

        # Extract stroke/strokeWidth
        stroke_info: dict = {}
        try:
            ln = shape.line
            if ln.width is not None:
                stroke_info["strokeWidth"] = round(_emu_to_px(ln.width))
            try:
                if ln.color.type == MSO_COLOR_TYPE.RGB:
                    stroke_info["stroke"] = f"#{ln.color.rgb}"
            except Exception:
                pass
        except Exception:
            pass

        rect_elem: dict = {
            "type": "rect",
            "x": x,
            "y": y,
            "width": w,
            "height": h,
            "fill": fill,
            **stroke_info,
        }
        elems.append(rect_elem)

        # If shape also has text, add a text element at same position
        try:
            if shape.has_text_frame:
                text_elem = self._parse_textbox(shape, x, y, w, h)
                if text_elem:
                    elems.append(text_elem)
        except Exception:
            pass

        return elems, warns

    def _parse_line(self, shape, x: int, y: int, w: int, h: int) -> Optional[dict]:
        """Parse a LINE shape into a line element dict."""
        stroke = "#000000"
        stroke_width = 1
        try:
            ln = shape.line
            if ln.color.type == MSO_COLOR_TYPE.RGB:
                stroke = f"#{ln.color.rgb}"
        except Exception:
            pass
        try:
            if shape.line.width is not None:
                stroke_width = max(1, round(_emu_to_px(shape.line.width)))
        except Exception:
            pass
        return {
            "type": "line",
            "x": x,
            "y": y,
            "width": w,
            "height": h,
            "stroke": stroke,
            "strokeWidth": stroke_width,
        }
