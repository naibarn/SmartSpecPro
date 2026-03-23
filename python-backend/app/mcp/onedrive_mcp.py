"""OneDrive MCP tool handlers.

Exposes OneDrive operations (search, read, list, info, read_excel) as tool
functions that are invoked via the internal MCP HTTP API.  Uses Microsoft
Graph API v1.0 via httpx (no SDK dependency).
"""

import logging
import math
import re
import time
from typing import Any, Optional
from urllib.parse import quote

import httpx

from app.services.mcp_client import _validate_mcp_url

logger = logging.getLogger(__name__)

GRAPH_BASE = "https://graph.microsoft.com/v1.0"

# Drive item ID validation: alphanumeric, hyphens, dots, exclamation marks, max 256 chars
_ITEM_ID_RE = re.compile(r"^[a-zA-Z0-9_.!-]{1,256}$")


# ── Exceptions ──────────────────────────────────────────────────────────────


class ToolError(Exception):
    """Raised for expected tool errors (file not found, permission denied, token expired)."""

    def __init__(self, code: str, message: str):
        self.code = code
        self.message = message
        super().__init__(message)


# ── Helpers ─────────────────────────────────────────────────────────────────


async def _get_access_token(user_id: int, token_service=None) -> str:
    """Obtain a valid Microsoft access token for the user."""
    if token_service is None:
        from app.services.microsoft_token_service import MicrosoftTokenService
        from app.core.database import get_db_context

        async with get_db_context() as db:
            svc = MicrosoftTokenService(db)
            return await svc.get_valid_access_token(user_id)
    return await token_service.get_valid_access_token(user_id)


def _validate_item_id(item_id: str) -> str:
    """Validate and return a safe drive item ID."""
    if not item_id or not _ITEM_ID_RE.match(item_id):
        raise ToolError("invalid_input", f"Invalid drive item ID format: {item_id!r}")
    return item_id


def _validate_query(query: str, max_length: int = 500) -> str:
    """Validate and sanitize a search query."""
    if not query or not query.strip():
        raise ToolError("invalid_input", "Search query cannot be empty")
    query = query.strip()
    if len(query) > max_length:
        raise ToolError("invalid_input", f"Search query too long (max {max_length} chars)")
    # Strip null bytes and control characters
    query = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", query)
    return query


def _handle_graph_error(status_code: int, body: str = ""):
    """Convert Microsoft Graph API errors to ToolError."""
    if status_code == 404:
        raise ToolError("file_not_found", "File not found or has been deleted.")
    if status_code == 403:
        raise ToolError("permission_denied", "You do not have permission to access this file.")
    if status_code == 401:
        raise ToolError("token_expired", "Microsoft account token has expired. Please reconnect.")


# ── Tool Handlers ───────────────────────────────────────────────────────────


async def search_onedrive_files(
    query: str,
    user_id: int,
    tenant_id: str,
    *,
    max_results: int = 10,
    token_service=None,
) -> dict[str, Any]:
    """Search the user's OneDrive for files matching a query.

    Free operation -- no credit charge.
    """
    from app.services.microsoft_token_service import InvalidGrantError

    query = _validate_query(query)
    max_results = max(1, min(max_results, 50))

    try:
        access_token = await _get_access_token(user_id, token_service)

        search_url = f"{GRAPH_BASE}/me/drive/root/search(q='{quote(query, safe='')}')"
        params = {
            "$select": "id,name,file,folder,size,lastModifiedDateTime,webUrl",
            "$top": str(max_results),
        }

        async with httpx.AsyncClient() as client:
            resp = await client.get(
                search_url,
                params=params,
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=15.0,
            )

        if resp.status_code != 200:
            _handle_graph_error(resp.status_code, resp.text)
            raise ToolError("api_error", f"OneDrive search failed (HTTP {resp.status_code})")

        data = resp.json()
        items = data.get("value", [])

        files = []
        for item in items:
            files.append({
                "id": item.get("id"),
                "name": item.get("name"),
                "mimeType": item.get("file", {}).get("mimeType", "folder" if "folder" in item else "unknown"),
                "size": item.get("size"),
                "lastModified": item.get("lastModifiedDateTime"),
                "webUrl": item.get("webUrl"),
            })

        return {"files": files, "total": len(files)}

    except InvalidGrantError:
        raise ToolError("token_expired", "Microsoft account token has expired. Please reconnect.")
    except ToolError:
        raise
    except Exception as e:
        logger.error("onedrive_search_error: %s", type(e).__name__)
        raise


async def read_onedrive_file(
    item_id: str,
    user_id: int,
    tenant_id: str,
    *,
    token_service=None,
    credit_charge_fn=None,
) -> dict[str, Any]:
    """Read the text content of a OneDrive file.

    Credits: max(1, ceil(char_count / 2000)), capped at 5.
    Service tag: onedrive.mcp_read
    """
    from app.services.microsoft_token_service import InvalidGrantError

    item_id = _validate_item_id(item_id)

    try:
        access_token = await _get_access_token(user_id, token_service)

        # Get file metadata
        meta_url = f"{GRAPH_BASE}/me/drive/items/{item_id}"
        async with httpx.AsyncClient() as client:
            meta_resp = await client.get(
                meta_url,
                params={"$select": "id,name,file,size"},
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=10.0,
            )

        if meta_resp.status_code != 200:
            _handle_graph_error(meta_resp.status_code, meta_resp.text)
            raise ToolError("api_error", f"Failed to get file metadata (HTTP {meta_resp.status_code})")

        meta = meta_resp.json()
        file_name = meta.get("name", "")
        mime_type = meta.get("file", {}).get("mimeType", "")

        # Download content — validate redirect targets to prevent SSRF
        download_url = f"{GRAPH_BASE}/me/drive/items/{item_id}/content"
        async with httpx.AsyncClient(follow_redirects=False) as client:
            dl_resp = await client.get(
                download_url,
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=30.0,
            )
            if dl_resp.status_code in (301, 302, 307, 308):
                location = dl_resp.headers.get("location", "")
                ssrf_err = _validate_mcp_url(location)
                if ssrf_err:
                    raise ToolError("ssrf_blocked", "Download redirect blocked: target URL failed SSRF validation")
                dl_resp = await client.get(location, headers={"Authorization": f"Bearer {access_token}"}, timeout=30.0)

        if dl_resp.status_code != 200:
            _handle_graph_error(dl_resp.status_code, dl_resp.text)
            raise ToolError("api_error", f"Failed to download file (HTTP {dl_resp.status_code})")

        # Extract text based on MIME type
        text = _extract_text(dl_resp.content, mime_type, file_name)

        if not text:
            return {"text": "", "file_name": file_name, "char_count": 0}

        # Calculate and charge credits
        char_count = len(text)
        credits = min(max(1, math.ceil(char_count / 2000)), 5)
        idempotency_key = f"mcp:read_onedrive_file:{item_id}:{int(time.time()) // 60}"

        if credit_charge_fn:
            await credit_charge_fn(
                user_id=user_id,
                amount=credits,
                service="onedrive.mcp_read",
                idempotency_key=idempotency_key,
                metadata={"item_id": item_id, "char_count": char_count},
            )
        else:
            from app.services.credit_billing_client import charge_credits_post_deduct

            await charge_credits_post_deduct(
                user_id=user_id,
                amount=credits,
                service="onedrive.mcp_read",
                idempotency_key=idempotency_key,
                metadata={"item_id": item_id, "char_count": char_count},
                source_type="indexing",
            )

        return {
            "text": text,
            "file_name": file_name,
            "char_count": char_count,
            "credits_charged": credits,
        }

    except InvalidGrantError:
        raise ToolError("token_expired", "Microsoft account token has expired. Please reconnect.")
    except ToolError:
        raise
    except Exception as e:
        logger.error("onedrive_read_error: %s", type(e).__name__)
        raise


async def read_excel_data(
    item_id: str,
    user_id: int,
    tenant_id: str,
    *,
    sheet_name: Optional[str] = None,
    cell_range: Optional[str] = None,
    token_service=None,
    credit_charge_fn=None,
) -> dict[str, Any]:
    """Read data from an Excel file via the Graph workbook API.

    Credits: max(1, ceil(total_cells / 500)), capped at 3.
    Service tag: onedrive.mcp_sheet
    """
    from app.services.microsoft_token_service import InvalidGrantError

    item_id = _validate_item_id(item_id)
    if cell_range and not re.match(r"^[A-Za-z0-9:!]+$", cell_range):
        raise ToolError("invalid_input", f"Invalid cell range format: {cell_range!r}")

    try:
        access_token = await _get_access_token(user_id, token_service)

        # Build range URL — URL-encode user-supplied worksheet/range to prevent injection
        worksheet = sheet_name or "Sheet1"
        worksheet_enc = quote(worksheet, safe="")
        if cell_range:
            cell_range_enc = quote(cell_range, safe="")
            range_url = f"{GRAPH_BASE}/me/drive/items/{item_id}/workbook/worksheets('{worksheet_enc}')/range(address='{cell_range_enc}')"
        else:
            range_url = f"{GRAPH_BASE}/me/drive/items/{item_id}/workbook/worksheets('{worksheet_enc}')/usedRange"

        async with httpx.AsyncClient() as client:
            resp = await client.get(
                range_url,
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=15.0,
            )

        if resp.status_code != 200:
            _handle_graph_error(resp.status_code, resp.text)
            raise ToolError("api_error", f"Failed to read Excel data (HTTP {resp.status_code})")

        data = resp.json()
        values = data.get("values", [])
        headers = values[0] if values else []
        rows = values[1:] if len(values) > 1 else []
        total_cells = sum(len(row) for row in values)

        # Calculate and charge credits
        credits = min(max(1, math.ceil(total_cells / 500)), 3)
        idempotency_key = f"mcp:read_excel_data:{item_id}:{int(time.time()) // 60}"

        if credit_charge_fn:
            await credit_charge_fn(
                user_id=user_id,
                amount=credits,
                service="onedrive.mcp_sheet",
                idempotency_key=idempotency_key,
                metadata={"item_id": item_id, "total_cells": total_cells},
            )
        else:
            from app.services.credit_billing_client import charge_credits_post_deduct

            await charge_credits_post_deduct(
                user_id=user_id,
                amount=credits,
                service="onedrive.mcp_sheet",
                idempotency_key=idempotency_key,
                metadata={"item_id": item_id, "total_cells": total_cells},
                source_type="indexing",
            )

        return {
            "sheet_name": worksheet,
            "headers": headers,
            "rows": rows,
            "total_cells": total_cells,
            "credits_charged": credits,
        }

    except InvalidGrantError:
        raise ToolError("token_expired", "Microsoft account token has expired. Please reconnect.")
    except ToolError:
        raise
    except Exception as e:
        logger.error("onedrive_excel_error: %s", type(e).__name__)
        raise


async def list_onedrive_folder(
    user_id: int,
    tenant_id: str,
    *,
    folder_id: Optional[str] = None,
    token_service=None,
) -> dict[str, Any]:
    """List files in a OneDrive folder.

    Free operation -- no credit charge.
    """
    from app.services.microsoft_token_service import InvalidGrantError

    if folder_id:
        folder_id = _validate_item_id(folder_id)

    try:
        access_token = await _get_access_token(user_id, token_service)

        if folder_id:
            list_url = f"{GRAPH_BASE}/me/drive/items/{folder_id}/children"
        else:
            list_url = f"{GRAPH_BASE}/me/drive/root/children"

        params = {
            "$select": "id,name,file,folder,size,lastModifiedDateTime",
            "$top": "100",
            "$orderby": "name",
        }

        async with httpx.AsyncClient() as client:
            resp = await client.get(
                list_url,
                params=params,
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=15.0,
            )

        if resp.status_code != 200:
            _handle_graph_error(resp.status_code, resp.text)
            raise ToolError("api_error", f"Failed to list folder (HTTP {resp.status_code})")

        data = resp.json()
        items = data.get("value", [])

        files = []
        for item in items:
            files.append({
                "id": item.get("id"),
                "name": item.get("name"),
                "mimeType": item.get("file", {}).get("mimeType", "folder" if "folder" in item else "unknown"),
                "size": item.get("size"),
                "lastModified": item.get("lastModifiedDateTime"),
                "isFolder": "folder" in item,
            })

        return {
            "files": files,
            "next_page_token": data.get("@odata.nextLink"),
        }

    except InvalidGrantError:
        raise ToolError("token_expired", "Microsoft account token has expired. Please reconnect.")
    except ToolError:
        raise
    except Exception as e:
        logger.error("onedrive_list_error: %s", type(e).__name__)
        raise


async def get_onedrive_file_info(
    item_id: str,
    user_id: int,
    tenant_id: str,
    *,
    token_service=None,
) -> dict[str, Any]:
    """Get metadata about a OneDrive file.

    Free operation -- no credit charge.
    """
    from app.services.microsoft_token_service import InvalidGrantError

    item_id = _validate_item_id(item_id)

    try:
        access_token = await _get_access_token(user_id, token_service)

        meta_url = f"{GRAPH_BASE}/me/drive/items/{item_id}"
        params = {"$select": "id,name,file,folder,size,lastModifiedDateTime,createdDateTime,webUrl,parentReference,createdBy,lastModifiedBy"}

        async with httpx.AsyncClient() as client:
            resp = await client.get(
                meta_url,
                params=params,
                headers={"Authorization": f"Bearer {access_token}"},
                timeout=10.0,
            )

        if resp.status_code != 200:
            _handle_graph_error(resp.status_code, resp.text)
            raise ToolError("api_error", f"Failed to get file info (HTTP {resp.status_code})")

        _SAFE_FILE_INFO_FIELDS = {"id", "name", "file", "folder", "size", "lastModifiedDateTime", "createdDateTime", "webUrl"}
        raw = resp.json()
        return {k: v for k, v in raw.items() if k in _SAFE_FILE_INFO_FIELDS}

    except InvalidGrantError:
        raise ToolError("token_expired", "Microsoft account token has expired. Please reconnect.")
    except ToolError:
        raise
    except Exception as e:
        logger.error("onedrive_info_error: %s", type(e).__name__)
        raise


# ── Text Extraction ─────────────────────────────────────────────────────────


def _extract_text(content: bytes, mime_type: str, file_name: str) -> str:
    """Extract plain text from file content based on MIME type."""
    mime_lower = mime_type.lower()

    # Plain text
    if "text/" in mime_lower or file_name.endswith((".txt", ".md", ".csv", ".json", ".xml", ".yaml", ".yml")):
        return content.decode("utf-8", errors="replace")

    # PDF
    if "pdf" in mime_lower or file_name.endswith(".pdf"):
        try:
            import io
            from PyPDF2 import PdfReader

            reader = PdfReader(io.BytesIO(content))
            texts = []
            for page in reader.pages:
                texts.append(page.extract_text() or "")
            return "\n".join(texts)
        except Exception as e:
            logger.warning("PDF extraction failed: %s", e)
            return ""

    # Word (.docx)
    if "wordprocessingml" in mime_lower or file_name.endswith(".docx"):
        try:
            import io
            from docx import Document

            doc = Document(io.BytesIO(content))
            return "\n".join(p.text for p in doc.paragraphs)
        except Exception as e:
            logger.warning("DOCX extraction failed: %s", e)
            return ""

    # PowerPoint (.pptx)
    if "presentationml" in mime_lower or file_name.endswith(".pptx"):
        try:
            import io
            from pptx import Presentation

            prs = Presentation(io.BytesIO(content))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
        except Exception as e:
            logger.warning("PPTX extraction failed: %s", e)
            return ""

    # Fallback: try as text
    try:
        return content.decode("utf-8", errors="strict")
    except (UnicodeDecodeError, ValueError):
        return ""


# ── Tool Registry ──────────────────────────────────────────────────────────

ONEDRIVE_TOOLS: list[dict[str, Any]] = [
    {
        "name": "search_onedrive_files",
        "description": "Search the user's OneDrive for files matching a query. Returns file names, types, and IDs.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Search query string"},
                "max_results": {
                    "type": "integer",
                    "default": 10,
                    "description": "Maximum results to return (1-50)",
                },
            },
            "required": ["query"],
        },
    },
    {
        "name": "read_onedrive_file",
        "description": "Read the text content of a OneDrive file. Supports Word, PowerPoint, PDF, and plain text files.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "item_id": {"type": "string", "description": "OneDrive drive item ID"},
            },
            "required": ["item_id"],
        },
    },
    {
        "name": "read_excel_data",
        "description": "Read data from an Excel file on OneDrive. Returns headers and rows. Optionally specify a sheet name and cell range.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "item_id": {"type": "string", "description": "OneDrive drive item ID for the Excel file"},
                "sheet_name": {"type": "string", "description": "Sheet tab name (default: Sheet1)"},
                "cell_range": {"type": "string", "description": "Cell range like A1:C10"},
            },
            "required": ["item_id"],
        },
    },
    {
        "name": "list_onedrive_folder",
        "description": "List files in a OneDrive folder. If no folder ID is provided, lists the root folder.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "folder_id": {"type": "string", "description": "Folder ID (omit for root)"},
            },
        },
    },
    {
        "name": "get_onedrive_file_info",
        "description": "Get metadata about a OneDrive file including name, type, size, modification date, and owner.",
        "inputSchema": {
            "type": "object",
            "properties": {
                "item_id": {"type": "string", "description": "OneDrive drive item ID"},
            },
            "required": ["item_id"],
        },
    },
]

# Handler dispatch map
TOOL_HANDLERS = {
    "search_onedrive_files": search_onedrive_files,
    "read_onedrive_file": read_onedrive_file,
    "read_excel_data": read_excel_data,
    "list_onedrive_folder": list_onedrive_folder,
    "get_onedrive_file_info": get_onedrive_file_info,
}
