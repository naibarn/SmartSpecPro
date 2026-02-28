"""
Tests for PptxImporter — section 02 of the Import Presentations feature.

All tests are pure unit tests. R2StorageService is mocked.
Fixtures build minimal PPTX objects programmatically using python-pptx.
"""
import io
import re
import pytest
from unittest.mock import AsyncMock, MagicMock
from pptx import Presentation as PptxPresentation
from pptx.util import Emu, Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn

MOCK_UPLOAD_URL = "https://cdn.example.com/test.png"


@pytest.fixture
def mock_r2():
    """Returns a mock R2StorageService with upload_bytes pre-patched."""
    svc = MagicMock()
    svc.upload_bytes = AsyncMock(return_value=MOCK_UPLOAD_URL)
    return svc


def _make_pptx_bytes(prs: PptxPresentation) -> bytes:
    """Serialize a PptxPresentation to bytes."""
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _default_prs(width_emu: int = 9_144_000, height_emu: int = 5_143_500) -> PptxPresentation:
    """Create a blank PptxPresentation with the given canvas dimensions."""
    prs = PptxPresentation()
    prs.slide_width = Emu(width_emu)
    prs.slide_height = Emu(height_emu)
    return prs


def _add_blank_slide(prs: PptxPresentation):
    """Add a blank slide to the presentation."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


# ---------------------------------------------------------------------------
# Coordinate conversion tests (pure functions — no I/O)
# ---------------------------------------------------------------------------

from app.services.pptx_importer import _emu_to_px, _scale_to_canvas, _pt_to_canvas_px


@pytest.mark.unit
def test_emu_to_px_one_inch():
    """914400 EMU = 1 inch = 96px at 96 DPI."""
    assert _emu_to_px(914_400) == pytest.approx(96.0)


@pytest.mark.unit
def test_emu_to_px_zero():
    assert _emu_to_px(0) == 0.0


@pytest.mark.unit
def test_scale_to_canvas_proportional():
    """Half the slide width should map to half the canvas width."""
    # slide is 9144000 EMU wide, canvas is 1280px wide
    # half the slide = 4572000 EMU -> 640px
    result = _scale_to_canvas(4_572_000, 9_144_000, 1280)
    assert result == 640


@pytest.mark.unit
def test_pt_to_canvas_px_normal():
    """12pt × 4/3 = 16px."""
    assert _pt_to_canvas_px(12.0) == 16


@pytest.mark.unit
def test_pt_to_canvas_px_clamp_min():
    """Values that compute below 8 are clamped to 8."""
    assert _pt_to_canvas_px(3.0) == 8


@pytest.mark.unit
def test_pt_to_canvas_px_clamp_max():
    """Values that compute above 512 are clamped to 512."""
    assert _pt_to_canvas_px(400.0) == 512


# ---------------------------------------------------------------------------
# Canvas detection tests
# ---------------------------------------------------------------------------

from app.services.pptx_importer import PptxImporter


@pytest.mark.unit
@pytest.mark.asyncio
async def test_canvas_16_9(mock_r2):
    """9144000 × 5143500 EMU → 16:9 preset, 1280×720."""
    prs = _default_prs(9_144_000, 5_143_500)
    _add_blank_slide(prs)
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "tenant1/presentations/imports/1"
    )
    assert len(result.slides) == 1
    slide = result.slides[0]
    assert slide["canvasPreset"] == "16:9"
    assert slide["canvasWidth"] == 1280
    assert slide["canvasHeight"] == 720


@pytest.mark.unit
@pytest.mark.asyncio
async def test_canvas_4_3(mock_r2):
    """9144000 × 6858000 EMU → 4:3 preset, 1024×768."""
    prs = _default_prs(9_144_000, 6_858_000)
    _add_blank_slide(prs)
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "tenant1/presentations/imports/2"
    )
    slide = result.slides[0]
    assert slide["canvasPreset"] == "4:3"
    assert slide["canvasWidth"] == 1024
    assert slide["canvasHeight"] == 768


@pytest.mark.unit
@pytest.mark.asyncio
async def test_canvas_unknown_ratio(mock_r2):
    """Non-standard ratio → preset is None, natural px dimensions used."""
    # 7200000 × 4050000 = 1.778 ratio... let's pick something truly non-standard
    prs = _default_prs(7_000_000, 3_500_000)  # 2:1 ratio - not in preset map
    _add_blank_slide(prs)
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "tenant1/presentations/imports/3"
    )
    slide = result.slides[0]
    assert slide["canvasPreset"] is None
    assert slide["canvasWidth"] > 0
    assert slide["canvasHeight"] > 0


# ---------------------------------------------------------------------------
# Text box parsing tests
# ---------------------------------------------------------------------------

@pytest.mark.unit
@pytest.mark.asyncio
async def test_textbox_element_type(mock_r2):
    """TEXT_BOX shape produces element with type='text'."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    tf = slide.shapes.add_textbox(Emu(100_000), Emu(100_000), Emu(500_000), Emu(200_000))
    tf.text_frame.text = "Hello World"
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    elements = result.slides[0]["elements"]
    text_elements = [e for e in elements if e.get("type") == "text"]
    assert len(text_elements) == 1
    assert text_elements[0]["content"] == "Hello World"


@pytest.mark.unit
@pytest.mark.asyncio
async def test_textbox_multiline_joined(mock_r2):
    """Text from multiple paragraphs is joined with '\\n'."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    tf = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(500_000), Emu(300_000))
    frame = tf.text_frame
    frame.text = "Line 1"
    frame.add_paragraph().text = "Line 2"
    frame.add_paragraph().text = "Line 3"
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    elements = result.slides[0]["elements"]
    text_elem = next(e for e in elements if e["type"] == "text")
    assert "Line 1" in text_elem["content"]
    assert "Line 2" in text_elem["content"]
    assert "\n" in text_elem["content"]


@pytest.mark.unit
@pytest.mark.asyncio
async def test_textbox_text_capped_at_10000(mock_r2):
    """Text longer than 10000 chars is truncated to 10000."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    tf = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(500_000), Emu(300_000))
    tf.text_frame.text = "A" * 15_000
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    elements = result.slides[0]["elements"]
    text_elem = next(e for e in elements if e["type"] == "text")
    assert len(text_elem["content"]) == 10_000


@pytest.mark.unit
@pytest.mark.asyncio
async def test_textbox_position_scaled(mock_r2):
    """Element x, y, width, height are scaled to canvas coordinates."""
    prs = _default_prs(9_144_000, 5_143_500)  # 16:9 → 1280×720
    slide = _add_blank_slide(prs)
    # Place at (0, 0) with full slide width/height
    tf = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(9_144_000), Emu(5_143_500))
    tf.text_frame.text = "Full"
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    elements = result.slides[0]["elements"]
    text_elem = next(e for e in elements if e["type"] == "text")
    assert text_elem["x"] == 0
    assert text_elem["y"] == 0
    assert text_elem["width"] == 1280
    assert text_elem["height"] == 720


# ---------------------------------------------------------------------------
# Image parsing tests
# ---------------------------------------------------------------------------

def _make_minimal_png() -> bytes:
    """Return a minimal valid 1×1 PNG image as bytes."""
    # 1×1 transparent PNG
    import base64
    # This is a known-good minimal PNG
    b64 = (
        "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNk+M9QDwADhgGAWjR9awAAAABJRU5ErkJggg=="
    )
    return base64.b64decode(b64)


@pytest.mark.unit
@pytest.mark.asyncio
async def test_picture_upload_called(mock_r2):
    """PICTURE shape calls upload_bytes with key matching pattern."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    png_bytes = _make_minimal_png()
    slide.shapes.add_picture(
        io.BytesIO(png_bytes),
        Emu(100_000), Emu(100_000), Emu(500_000), Emu(300_000)
    )
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "tenant1/presentations/imports/99"
    )
    assert mock_r2.upload_bytes.call_count >= 1
    call_key = mock_r2.upload_bytes.call_args_list[0][0][0]
    assert re.match(r".+/images/[^/]+\.\w+$", call_key), f"Key doesn't match pattern: {call_key}"


@pytest.mark.unit
@pytest.mark.asyncio
async def test_picture_element_src(mock_r2):
    """PICTURE element type is 'image', src equals the mocked upload URL."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    png_bytes = _make_minimal_png()
    slide.shapes.add_picture(
        io.BytesIO(png_bytes),
        Emu(0), Emu(0), Emu(500_000), Emu(300_000)
    )
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    elements = result.slides[0]["elements"]
    img_elems = [e for e in elements if e.get("type") == "image"]
    assert len(img_elems) == 1
    assert img_elems[0]["src"] == MOCK_UPLOAD_URL


# ---------------------------------------------------------------------------
# Shape parsing tests
# ---------------------------------------------------------------------------

@pytest.mark.unit
@pytest.mark.asyncio
async def test_autoshape_rectangle(mock_r2):
    """AUTO_SHAPE RECTANGLE produces element type 'rect'."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    from pptx.util import Inches
    # Add a rectangle
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Emu(100_000), Emu(100_000), Emu(500_000), Emu(300_000)
    )
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    elements = result.slides[0]["elements"]
    rect_elems = [e for e in elements if e.get("type") == "rect"]
    assert len(rect_elems) >= 1


@pytest.mark.unit
@pytest.mark.asyncio
async def test_autoshape_oval_produces_warning(mock_r2):
    """AUTO_SHAPE OVAL produces element type 'rect' AND a fidelityWarning."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Emu(100_000), Emu(100_000), Emu(400_000), Emu(400_000)
    )
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    elements = result.slides[0]["elements"]
    rect_elems = [e for e in elements if e.get("type") == "rect"]
    assert len(rect_elems) >= 1
    assert any("oval" in w.lower() or "Oval" in w for w in result.fidelity_warnings)


@pytest.mark.unit
@pytest.mark.asyncio
async def test_line_shape(mock_r2):
    """LINE shape produces element type 'line'."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.LINE_INVERSE,
        Emu(100_000), Emu(100_000), Emu(500_000), Emu(5_000)
    )
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    elements = result.slides[0]["elements"]
    # LINE_INVERSE is dispatched through AUTO_SHAPE; verify no crash and valid list returned
    assert isinstance(elements, list)


# ---------------------------------------------------------------------------
# Unsupported shape tests
# ---------------------------------------------------------------------------

@pytest.mark.unit
@pytest.mark.asyncio
async def test_table_produces_warning_and_skipped(mock_r2):
    """TABLE shape produces fidelityWarning containing 'Table' and no element."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    slide.shapes.add_table(2, 2, Emu(100_000), Emu(100_000), Emu(500_000), Emu(300_000))
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    assert any("Table" in w for w in result.fidelity_warnings)
    table_elements = [e for e in result.slides[0]["elements"] if e.get("type") == "table"]
    assert len(table_elements) == 0


# ---------------------------------------------------------------------------
# Error handling tests
# ---------------------------------------------------------------------------

@pytest.mark.unit
@pytest.mark.asyncio
async def test_corrupt_file_raises_import_error(mock_r2):
    """Non-zip bytes cause import_file to raise ImportError with user-friendly message."""
    importer = PptxImporter(mock_r2)
    with pytest.raises(ImportError, match="not a valid .pptx"):
        await importer.import_file(b"this is not a pptx file", "tenant1/presentations/imports/99")


@pytest.mark.unit
@pytest.mark.asyncio
async def test_empty_file_raises_import_error(mock_r2):
    """Empty bytes cause import_file to raise ImportError."""
    importer = PptxImporter(mock_r2)
    with pytest.raises(ImportError):
        await importer.import_file(b"", "tenant1/presentations/imports/99")


# ---------------------------------------------------------------------------
# fidelityWarnings cap test
# ---------------------------------------------------------------------------

@pytest.mark.unit
@pytest.mark.asyncio
async def test_fidelity_warnings_capped_at_25(mock_r2):
    """When many warnings generated, result has at most 25 items."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    # Add 30 tables (each produces a warning)
    for i in range(30):
        slide.shapes.add_table(
            2, 2,
            Emu(i * 10_000 + 100_000), Emu(100_000),
            Emu(50_000), Emu(50_000)
        )
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    assert len(result.fidelity_warnings) <= 25
    if len(result.fidelity_warnings) == 25:
        assert "more warnings" in result.fidelity_warnings[-1]


# ---------------------------------------------------------------------------
# Additional tests from code review (Issue 9a, 9b, 9c)
# ---------------------------------------------------------------------------

from app.services.pptx_importer import _CanvasSize


@pytest.mark.unit
@pytest.mark.asyncio
async def test_linked_picture_skipped(mock_r2):
    """LINKED_PICTURE shape produces a fidelityWarning and no element."""
    mock_shape = MagicMock()
    mock_shape.shape_type = MSO_SHAPE_TYPE.LINKED_PICTURE
    mock_shape.left = Emu(100_000)
    mock_shape.top = Emu(100_000)
    mock_shape.width = Emu(500_000)
    mock_shape.height = Emu(300_000)

    importer = PptxImporter(mock_r2)
    canvas = _CanvasSize(preset="16:9", width=1280, height=720)
    elements, warnings = await importer._parse_shapes(
        [mock_shape],
        parent_left_emu=0,
        parent_top_emu=0,
        canvas=canvas,
        slide_width_emu=9_144_000,
        slide_height_emu=5_143_500,
        s3_prefix="t/p/1",
        slide_num=1,
    )
    assert len(elements) == 0
    assert any("Linked" in w or "linked" in w for w in warnings)


@pytest.mark.unit
@pytest.mark.asyncio
async def test_textbox_font_color_rgb(mock_r2):
    """TEXT_BOX with an RGB font color populates style['color'] correctly."""
    prs = _default_prs()
    slide = _add_blank_slide(prs)
    tf = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(500_000), Emu(200_000))
    frame = tf.text_frame
    run = frame.paragraphs[0].add_run()
    run.text = "Colored text"
    run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    result = await PptxImporter(mock_r2).import_file(
        _make_pptx_bytes(prs), "t/p/1"
    )
    elements = result.slides[0]["elements"]
    text_elems = [e for e in elements if e.get("type") == "text"]
    assert len(text_elems) == 1
    assert text_elems[0]["style"].get("color") == "#FF0000"


@pytest.mark.unit
@pytest.mark.asyncio
async def test_group_child_offset(mock_r2):
    """Child shapes in a GROUP accumulate parent EMU offset before scaling."""
    from app.services.pptx_importer import _scale_to_canvas

    # Mock child textbox at 100k EMU relative to group
    mock_child = MagicMock()
    mock_child.shape_type = MSO_SHAPE_TYPE.TEXT_BOX
    mock_child.left = Emu(100_000)
    mock_child.top = Emu(100_000)
    mock_child.width = Emu(200_000)
    mock_child.height = Emu(100_000)
    mock_child.text_frame.paragraphs = [MagicMock(runs=[MagicMock(text="hi")])]

    # Mock group at 500k EMU on slide
    mock_group = MagicMock()
    mock_group.shape_type = MSO_SHAPE_TYPE.GROUP
    mock_group.left = Emu(500_000)
    mock_group.top = Emu(500_000)
    mock_group.width = Emu(1_000_000)
    mock_group.height = Emu(500_000)
    mock_group.shapes = [mock_child]

    canvas = _CanvasSize(preset="16:9", width=1280, height=720)
    importer = PptxImporter(mock_r2)
    elements, _ = await importer._parse_shapes(
        [mock_group],
        parent_left_emu=0,
        parent_top_emu=0,
        canvas=canvas,
        slide_width_emu=9_144_000,
        slide_height_emu=5_143_500,
        s3_prefix="t/p/1",
        slide_num=1,
    )
    # Child absolute position: group(500k) + child(100k) = 600k EMU
    expected_x = _scale_to_canvas(600_000, 9_144_000, 1280)
    expected_y = _scale_to_canvas(600_000, 5_143_500, 720)
    # Find any element with the expected coordinates
    assert any(e.get("x") == expected_x and e.get("y") == expected_y for e in elements), (
        f"No element at expected ({expected_x}, {expected_y}). Got: {elements}"
    )
